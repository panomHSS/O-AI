from pathlib import Path

from pptx import Presentation

from app.readers.base import DocumentExtractionError, SourceSection


class PptxDocumentReader:
    extensions = frozenset({".pptx"})

    def extract(self, path: Path) -> list[SourceSection]:
        try:
            presentation = Presentation(path)
            sections: list[SourceSection] = []
            for index, slide in enumerate(presentation.slides, start=1):
                values = [shape.text.strip() for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.text.strip()]
                if values:
                    sections.append(SourceSection(index, "\n".join(values), f"PowerPoint slide {index}", {"atomic": "true"}))
        except Exception as error:
            raise DocumentExtractionError("Unable to extract text from this PowerPoint presentation.") from error

        if not sections:
            raise DocumentExtractionError("This PowerPoint presentation has no extractable text.")
        return sections
