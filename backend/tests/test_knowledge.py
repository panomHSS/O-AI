import asyncio
import tempfile
import unittest
from pathlib import Path
from uuid import UUID, uuid4

from docx import Document as DocxFile
from openpyxl import Workbook
from pypdf import PdfWriter
from pptx import Presentation
from sqlalchemy.orm import sessionmaker

from app.api.dependencies import get_knowledge_service
from app.db.session import create_database_engine, get_db, initialize_database
from app.readers import create_document_reader_registry
from app.readers.base import DocumentExtractionError
from app.repositories.knowledge import KnowledgeRepository
from app.services.knowledge import KnowledgeScanConflictError, KnowledgeSearchValidationError, KnowledgeService
from app.main import app
from test_api_standardization import invoke_app


class KnowledgeTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary_directory = tempfile.TemporaryDirectory()
        self.root = Path(self.temporary_directory.name) / "knowledge"
        self.root.mkdir()
        self.engine = create_database_engine(f"sqlite:///{(Path(self.temporary_directory.name) / 'oai.db').as_posix()}")
        initialize_database(self.engine)
        self.Session = sessionmaker(bind=self.engine, autoflush=False, autocommit=False, expire_on_commit=False)
        self.sessions = []

        def test_database_session():
            session = self.Session()
            try:
                yield session
            finally:
                session.close()

        app.dependency_overrides[get_db] = test_database_session
        def test_knowledge_service():
            service = self._service()
            try:
                yield service
            finally:
                service._repository.close()

        app.dependency_overrides[get_knowledge_service] = test_knowledge_service

    def tearDown(self) -> None:
        app.dependency_overrides.clear()
        for session in self.sessions:
            session.close()
        self.engine.dispose()
        self.temporary_directory.cleanup()

    def _service(self) -> KnowledgeService:
        session = self.Session()
        self.sessions.append(session)
        return KnowledgeService(KnowledgeRepository(session), create_document_reader_registry(), str(self.root), 1, 40, 10)

    def request(self, *args, **kwargs):
        return asyncio.run(invoke_app(*args, **kwargs))

    def scan(self):
        return self.request("/api/v1/knowledge/scan", method="POST")

    def test_readers_extract_supported_formats(self) -> None:
        (self.root / "note.txt").write_text("plain text", encoding="utf-8")
        (self.root / "note.md").write_text("# markdown", encoding="utf-8")
        (self.root / "page.html").write_text("<h1>HTML text</h1><script>ignored()</script>", encoding="utf-8")
        (self.root / "mail.eml").write_text("From: sender@example.test\nSubject: Hello\nContent-Type: text/plain; charset=utf-8\n\nEmail body", encoding="utf-8")
        (self.root / "rows.csv").write_text("name,value\nalpha,1\n", encoding="utf-8")
        docx = DocxFile(); docx.add_paragraph("Word text"); docx.save(self.root / "word.docx")
        workbook = Workbook(); workbook.active.append(["Excel", "value"]); workbook.save(self.root / "sheet.xlsx")
        presentation = Presentation(); presentation.slides.add_slide(presentation.slide_layouts[6]).shapes.add_textbox(0, 0, 100, 100).text_frame.text = "Slide text"; presentation.save(self.root / "slides.pptx")
        writer = PdfWriter(); writer.add_blank_page(width=72, height=72); writer.write(self.root / "empty.pdf")

        status, _, body = self.scan()
        self.assertEqual(status, 200)
        self.assertEqual(body["data"]["indexed"], 8)
        self.assertEqual(body["data"]["failed"], 1)

    def test_unsupported_and_maximum_size_are_reported(self) -> None:
        (self.root / "unknown.bin").write_bytes(b"x")
        (self.root / "large.txt").write_text("x" * (1024 * 1024 + 1), encoding="utf-8")
        _, _, body = self.scan()
        self.assertEqual(body["data"]["unsupported"], 1)
        self.assertEqual(body["data"]["failed"], 1)

    def test_duplicate_content_at_different_paths_remains_separate(self) -> None:
        (self.root / "first.txt").write_text("same source text", encoding="utf-8")
        nested = self.root / "nested"; nested.mkdir(); (nested / "second.txt").write_text("same source text", encoding="utf-8")
        self.scan()
        status, _, body = self.request("/api/v1/knowledge/documents")
        self.assertEqual(status, 200)
        self.assertEqual(body["data"]["total"], 2)
        self.assertEqual({item["source_path"] for item in body["data"]["items"]}, {"first.txt", "nested/second.txt"})
        self.assertNotIn(str(self.root), str(body))

    def test_discovery_prevents_path_escape_and_symlink_escape(self) -> None:
        outside = Path(self.temporary_directory.name) / "outside.txt"
        outside.write_text("outside", encoding="utf-8")
        (self.root / "inside.txt").write_text("inside", encoding="utf-8")
        self.assertEqual([path.relative_to(self.root).as_posix() for path in self._service()._discover(self.root)], ["inside.txt"])
        try:
            (self.root / "escape.txt").symlink_to(outside)
        except OSError:
            self.skipTest("Symlinks are not available in this Windows test environment.")
        _, _, listing = self.scan()
        self.assertEqual({item["source_path"] for item in self.request("/api/v1/knowledge/documents")[2]["data"]["items"]}, {"inside.txt"})

    def test_unchanged_changed_and_failed_reindex_preserves_old_chunks(self) -> None:
        source = self.root / "note.txt"; source.write_text("first indexed words", encoding="utf-8")
        self.scan()
        _, _, unchanged = self.scan()
        self.assertEqual(unchanged["data"]["unchanged"], 1)
        source.write_text("second indexed words", encoding="utf-8")
        _, _, changed = self.scan()
        self.assertEqual(changed["data"]["indexed"], 1)
        source.write_text("", encoding="utf-8")
        _, _, failed = self.scan()
        self.assertEqual(failed["data"]["failed"], 1)
        _, _, results = self.request("/api/v1/knowledge/search?q=second")
        self.assertEqual(results["data"]["items"][0]["source_path"], "note.txt")

    def test_chunks_are_ordered(self) -> None:
        (self.root / "long.txt").write_text("one two three four five six seven eight nine ten eleven twelve", encoding="utf-8")
        self.scan()
        _, _, listing = self.request("/api/v1/knowledge/documents")
        _, _, detail = self.request(f"/api/v1/knowledge/documents/{listing['data']['items'][0]['id']}")
        indexes = [chunk["chunk_index"] for chunk in detail["data"]["chunks"]]
        self.assertEqual(indexes, list(range(len(indexes))))

    def test_missing_is_retained_but_excluded_from_search(self) -> None:
        source = self.root / "note.txt"; source.write_text("find me", encoding="utf-8")
        self.scan(); source.unlink(); self.scan()
        _, _, listing = self.request("/api/v1/knowledge/documents?status=missing")
        self.assertEqual(listing["data"]["total"], 1)
        _, _, search = self.request("/api/v1/knowledge/search?q=find")
        self.assertEqual(search["data"]["items"], [])

    def test_document_api_pagination_detail_delete_and_request_id(self) -> None:
        (self.root / "one.txt").write_text("one", encoding="utf-8")
        (self.root / "two.txt").write_text("two", encoding="utf-8")
        self.scan()
        status, headers, listing = self.request("/api/v1/knowledge/documents?page=1&page_size=1", headers={"X-Request-ID": "knowledge-request"})
        self.assertEqual(status, 200); self.assertEqual(headers["x-request-id"], "knowledge-request"); self.assertEqual(len(listing["data"]["items"]), 1)
        document_id = listing["data"]["items"][0]["id"]; UUID(document_id)
        status, _, detail = self.request(f"/api/v1/knowledge/documents/{document_id}")
        self.assertEqual(status, 200); self.assertGreater(detail["data"]["chunk_count"], 0)
        source = self.root / detail["data"]["source_path"]
        status, _, _ = self.request(f"/api/v1/knowledge/documents/{document_id}", method="DELETE")
        self.assertEqual(status, 200); self.assertTrue(source.exists())

    def test_invalid_id_and_malformed_search_are_safe_envelopes(self) -> None:
        status, _, invalid = self.request("/api/v1/knowledge/documents/not-a-uuid")
        self.assertEqual(status, 422); self.assertEqual(invalid["error"]["code"], "VALIDATION_ERROR")
        status, _, malformed = self.request("/api/v1/knowledge/search?q=!!!")
        self.assertEqual(status, 422); self.assertEqual(malformed["error"]["code"], "KNOWLEDGE_SEARCH_VALIDATION_ERROR")

    def test_interrupted_scan_does_not_mark_missing(self) -> None:
        (self.root / "note.txt").write_text("survive", encoding="utf-8")
        service = self._service(); service.scan()
        original_discover = service._discover
        def interrupted(root):
            yield from original_discover(root)
            raise RuntimeError("interrupted")
        service._discover = interrupted  # type: ignore[method-assign]
        with self.assertRaises(RuntimeError): service.scan()
        self.assertEqual(service.list_documents(1, 10, "missing").total, 0)

    def test_concurrent_scan_conflict_and_chunk_configuration(self) -> None:
        service = self._service()
        self.assertTrue(service._scan_lock.acquire(blocking=False))
        try:
            with self.assertRaises(KnowledgeScanConflictError): service.scan()
        finally:
            service._scan_lock.release()
        with self.assertRaises(ValueError):
            from app.core.config import Settings
            Settings(oai_chunk_size_chars=10, oai_chunk_overlap_chars=10)

    def test_reader_rejects_empty_pdf(self) -> None:
        path = self.root / "empty.pdf"; writer = PdfWriter(); writer.add_blank_page(width=72, height=72); writer.write(path)
        with self.assertRaises(DocumentExtractionError):
            create_document_reader_registry().get(path).extract(path)  # type: ignore[union-attr]
